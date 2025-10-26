import os
import tempfile
import uuid
import asyncio
import aiohttp
import aiofiles
import random
from typing import Optional, Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import selectinload

from database_models import Presentation as PresentationModel


class PPTXService:
    
    class FilteredSlideData:
        """Wrapper for slide data with filtered content"""
        def __init__(self, title, content, image_url, image_alt, layout):
            self.title = title
            self.content = content
            self.image_url = image_url
            self.image_alt = image_alt
            self.layout = layout
    
    @staticmethod
    def _filter_content_for_layout(content: str, layout_type: str) -> str:
        """Filter content based on layout type - remove non-bullet lines for certain layouts"""
        if not content:
            return content
        
        # Layouts that should only show bullet points (no key insights)
        bullet_only_layouts = ['image_left', 'image_right', 'image_top']
        
        if layout_type in bullet_only_layouts:
            # Keep only lines that start with bullet points
            lines = content.split('\n')
            bullet_lines = [line for line in lines if line.strip().startswith('-')]
            return '\n'.join(bullet_lines)
        
        # For other layouts, return all content
        return content
    
    # Theme color mappings matching frontend
    THEMES = {
        'minimal': {
            'title_bg': RGBColor(99, 102, 241),    # indigo-600
            'accent': RGBColor(99, 102, 241),       # indigo-500
            'circle1': RGBColor(99, 102, 241),      # indigo-600
            'circle2': RGBColor(147, 51, 234),      # purple-600
        },
        'professional': {
            'title_bg': RGBColor(30, 41, 59),       # slate-800
            'accent': RGBColor(37, 99, 235),        # blue-600
            'circle1': RGBColor(37, 99, 235),       # blue-600
            'circle2': RGBColor(59, 130, 246),      # blue-500
        },
        'creative': {
            'title_bg': RGBColor(147, 51, 234),     # purple-600
            'accent': RGBColor(236, 72, 153),       # pink-500
            'circle1': RGBColor(147, 51, 234),      # purple-600
            'circle2': RGBColor(236, 72, 153),      # pink-500
        },
        'academic': {
            'title_bg': RGBColor(4, 120, 87),       # emerald-700
            'accent': RGBColor(5, 150, 105),        # emerald-600
            'circle1': RGBColor(5, 150, 105),       # emerald-600
            'circle2': RGBColor(16, 185, 129),      # emerald-500
        },
        'dark': {
            'title_bg': RGBColor(17, 24, 39),       # gray-900
            'content_bg': RGBColor(17, 24, 39),     # gray-900 for content slides too
            'text_color': RGBColor(243, 244, 246),  # gray-100 for text
            'accent': RGBColor(34, 211, 238),       # cyan-400
            'circle1': RGBColor(34, 211, 238),      # cyan-400
            'circle2': RGBColor(6, 182, 212),       # cyan-500
        },
    }
    
    def __init__(self):
        pass
    
    def _get_theme_colors(self, style: str) -> dict:
        """Get theme colors based on presentation style"""
        return self.THEMES.get(style, self.THEMES['minimal'])
    
    async def export_to_pptx(self, db: AsyncSession, presentation_id: uuid.UUID) -> str:
        """Export presentation to PPTX and return file path"""
        
        # Get presentation with slides
        result = await db.execute(
            select(PresentationModel)
            .options(selectinload(PresentationModel.slides))
            .where(PresentationModel.id == presentation_id)
        )
        
        presentation_data = result.scalar_one_or_none()
        if not presentation_data:
            raise ValueError("Презентация не найдена")
        
        # Sort slides by slide_number
        slides = sorted(presentation_data.slides, key=lambda s: s.slide_number)
        
        # Create temporary directory for downloading images
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Download all images first
            image_paths = await self._download_images(slides, temp_dir)
            
            # Create PPTX presentation
            ppt = Presentation()
            
            # Set slide size to 16:9
            ppt.slide_width = Inches(13.33)
            ppt.slide_height = Inches(7.5)
            
            # Get layout order from presentation (fallback to None for old presentations)
            layout_order = getattr(presentation_data, 'layout_order', None)
            
            for slide_data in slides:
                self._add_slide(ppt, slide_data, image_paths.get(slide_data.slide_number), layout_order, presentation_data.style)
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                ppt.save(temp_file.name)
                return temp_file.name
                
        finally:
            # Clean up temp directory
            try:
                import shutil
                shutil.rmtree(temp_dir, ignore_errors=True)
            except:
                pass
    
    async def _download_images(self, slides, temp_dir: str) -> Dict[int, str]:
        """Download all images and return mapping of slide_number to local path"""
        image_paths = {}
        
        download_tasks = []
        slide_numbers = []
        
        for slide in slides:
            if slide.image_url:
                slide_numbers.append(slide.slide_number)
                download_tasks.append(self._download_image(slide.image_url, temp_dir))
        
        if download_tasks:
            downloaded_paths = await asyncio.gather(*download_tasks, return_exceptions=True)
            
            for slide_number, path in zip(slide_numbers, downloaded_paths):
                if isinstance(path, str) and os.path.exists(path):
                    image_paths[slide_number] = path
                    
        return image_paths
    
    async def _download_image(self, image_url: str, temp_dir: str) -> Optional[str]:
        """Download a single image and return local path"""
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(image_url) as response:
                    if response.status == 200:
                        # Generate unique filename
                        file_extension = '.jpg'  # Default to jpg
                        content_type = response.headers.get('Content-Type', '')
                        if 'png' in content_type:
                            file_extension = '.png'
                        elif 'webp' in content_type:
                            file_extension = '.webp'
                            
                        filename = f"{uuid.uuid4()}{file_extension}"
                        file_path = os.path.join(temp_dir, filename)
                        
                        # Download and save file
                        async with aiofiles.open(file_path, 'wb') as f:
                            async for chunk in response.content.iter_chunked(8192):
                                await f.write(chunk)
                        
                        # Convert to JPG if needed for better PPTX compatibility
                        if file_extension != '.jpg':
                            jpg_path = await self._convert_to_jpg(file_path, temp_dir)
                            if jpg_path:
                                return jpg_path
                        
                        return file_path
        except Exception as e:
            print(f"Error downloading image {image_url}: {e}")
            return None
    
    async def _convert_to_jpg(self, image_path: str, temp_dir: str) -> Optional[str]:
        """Convert image to JPG format"""
        try:
            with Image.open(image_path) as img:
                # Convert RGBA to RGB for JPG
                if img.mode in ('RGBA', 'LA'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    img = background
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                jpg_path = os.path.join(temp_dir, f"{uuid.uuid4()}.jpg")
                img.save(jpg_path, 'JPEG', quality=90)
                
                # Remove original file
                os.remove(image_path)
                
                return jpg_path
        except Exception as e:
            print(f"Error converting image to JPG: {e}")
            return image_path  # Return original on error
    
    def _add_slide(self, ppt: Presentation, slide_data, image_path: Optional[str], layout_order: Optional[List[str]] = None, presentation_style: str = 'minimal'):
        """Add a slide to the presentation"""
        
        # Use blank slide layout
        blank_slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(blank_slide_layout)
        
        if slide_data.layout == 'title-slide':
            self._create_title_slide(slide, slide_data, presentation_style)
        else:
            # Add variety to slides based on slide number and layout order
            slide_style = self._get_slide_style(slide_data.slide_number, layout_order, presentation_style)
            self._create_content_slide(slide, slide_data, image_path, slide_style)
    
    def _create_title_slide(self, slide, slide_data, presentation_style: str = 'minimal'):
        """Create a title slide with gradient background"""
        
        theme_colors = self._get_theme_colors(presentation_style)
        
        # Set background to theme color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme_colors['title_bg']
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.33), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        # Title font
        title_font = title_para.font
        title_font.name = 'Inter'
        title_font.size = Pt(54)
        title_font.bold = True
        title_font.color.rgb = RGBColor(255, 255, 255)
        
        # Add content if exists
        if slide_data.content:
            content_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.5), Inches(9.33), Inches(1.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_para = content_frame.paragraphs[0]
            content_para.text = self._clean_markdown(slide_data.content)
            content_para.alignment = PP_ALIGN.CENTER
            
            # Content font
            content_font = content_para.font
            content_font.name = 'Inter'
            content_font.size = Pt(20)
            content_font.color.rgb = RGBColor(255, 255, 255)
    
    def _get_slide_style(self, slide_number: int, layout_order: Optional[List[str]] = None, presentation_style: str = 'minimal') -> dict:
        """Get style variations for different slides to avoid monotony"""
        
        theme_colors = self._get_theme_colors(presentation_style)
        
        # Common theme properties for all slides
        common_props = {
            "content_bg": theme_colors.get('content_bg', RGBColor(255, 255, 255)),
            "text_color": theme_colors.get('text_color', RGBColor(17, 24, 39)),
        }
        
        styles = [
            {
                "layout": "image_left",
                "bullet_type": "numbers",
                "accent_color": theme_colors['accent'],
                "circle_colors": [theme_colors['circle1'], theme_colors['circle2']],
                "text_alignment": "left",
                **common_props
            },
            {
                "layout": "image_right", 
                "bullet_type": "dots", 
                "accent_color": theme_colors['accent'],
                "circle_colors": [theme_colors['circle1'], theme_colors['circle2']],
                "text_alignment": "left",
                **common_props
            },
            {
                "layout": "text_only",
                "bullet_type": "numbers",
                "accent_color": theme_colors['accent'],
                "circle_colors": [theme_colors['circle1'], theme_colors['circle2']],
                "text_alignment": "center",
                **common_props
            },
            {
                "layout": "split_content",
                "bullet_type": "dots",
                "accent_color": theme_colors['accent'],
                "circle_colors": [theme_colors['circle1'], theme_colors['circle2']],
                "text_alignment": "left",
                **common_props
            },
            {
                "layout": "image_top",
                "bullet_type": "numbers", 
                "accent_color": theme_colors['accent'],
                "circle_colors": [theme_colors['circle1'], theme_colors['circle2']],
                "text_alignment": "center",
                **common_props
            },
            {
                "layout": "grid_layout",
                "bullet_type": "dots",
                "accent_color": theme_colors['accent'],
                "circle_colors": [theme_colors['circle1'], theme_colors['circle2']],
                "text_alignment": "left",
                **common_props
            }
        ]
        
        # For content slides (not title), adjust index to match frontend
        # Slide 1 is title (not using this method), slide 2+ are content
        # So slide 2 should use index 0, slide 3 should use index 1, etc.
        content_slide_index = slide_number - 2  # -2 because slide 1 is title
        
        # If layout_order provided, use it; otherwise use sequential
        if layout_order and content_slide_index >= 0:
            layout_name = layout_order[content_slide_index % len(layout_order)]
            # Find the matching style
            for style in styles:
                if style["layout"] == layout_name:
                    return style
        
        # Fallback for old presentations or edge cases
        return styles[max(0, content_slide_index) % len(styles)]
    
    def _create_content_slide(self, slide, slide_data, image_path: Optional[str], slide_style: dict):
        """Create a content slide with image and text using different layout styles"""
        
        # Set background color based on theme
        background = slide.background
        fill = background.fill
        fill.solid()
        # Use theme background if available (for dark mode), otherwise white
        bg_color = slide_style.get('content_bg', RGBColor(255, 255, 255))
        fill.fore_color.rgb = bg_color
        
        # Route to different layout methods based on style
        layout_type = slide_style.get("layout", "image_left")
        
        if layout_type == "image_left":
            self._create_image_left_layout(slide, slide_data, image_path, slide_style)
        elif layout_type == "image_right":
            self._create_image_right_layout(slide, slide_data, image_path, slide_style)
        elif layout_type == "text_only":
            self._create_text_only_layout(slide, slide_data, slide_style)
        elif layout_type == "split_content":
            self._create_split_content_layout(slide, slide_data, image_path, slide_style)
        elif layout_type == "image_top":
            self._create_image_top_layout(slide, slide_data, image_path, slide_style)
        elif layout_type == "grid_layout":
            self._create_grid_layout(slide, slide_data, image_path, slide_style)
        else:
            # Fallback to image_left
            self._create_image_left_layout(slide, slide_data, image_path, slide_style)
    
    def _create_image_left_layout(self, slide, slide_data, image_path: Optional[str], slide_style: dict):
        """Original layout - image on left, text on right"""
        
        # Filter content to show only bullet points
        filtered_slide_data = self.FilteredSlideData(
            title=slide_data.title,
            content=self._filter_content_for_layout(slide_data.content, 'image_left'),
            image_url=slide_data.image_url,
            image_alt=slide_data.image_alt,
            layout=slide_data.layout
        )
        
        if image_path and os.path.exists(image_path):
            # Add image on the left with styling to match web UI
            try:
                # Add white background container (mimicking the web UI design)
                bg_container = slide.shapes.add_shape(
                    1,  # Rectangle shape
                    Inches(0.4), Inches(1.4), Inches(5.5), Inches(4.2)  # Reduced width from 6.2 to 5.5
                )
                bg_container.fill.solid()
                bg_container.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                bg_container.line.color.rgb = RGBColor(229, 231, 235)  # Light gray border
                bg_container.line.width = Pt(1)
                
                # Add rounded corners (limited support in PPTX)
                try:
                    # Add rounded rectangle effect
                    from pptx.oxml.xmlchemy import OxmlElement
                    from lxml import etree
                    sp_element = bg_container._element
                    sp_pr = sp_element.xpath("p:spPr")[0]
                    nsmap = sp_pr.nsmap
                    
                    # Add rounded rectangle geometry
                    prst_geom = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 8333")  # Smaller rounded corner radius
                except:
                    pass  # Rounded corners are optional
                
                # Add shadow effect (simplified)
                try:
                    from pptx.oxml.xmlchemy import OxmlElement
                    from lxml import etree
                    sp_element = bg_container._element
                    sp_pr = sp_element.xpath("p:spPr")[0]
                    nsmap = sp_pr.nsmap
                    effect_list = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}effectLst", nsmap=nsmap)
                    outer_shadow = etree.SubElement(
                        effect_list, f"{{{nsmap['a']}}}outerShdw",
                        {"blurRad": "38100", "dist": "25400", "dir": "5400000", "algn": "tl"},
                        nsmap=nsmap
                    )
                    color_element = etree.SubElement(
                        outer_shadow, f"{{{nsmap['a']}}}srgbClr", {"val": "000000"}, nsmap=nsmap
                    )
                    etree.SubElement(
                        color_element, f"{{{nsmap['a']}}}alpha", {"val": "15000"}, nsmap=nsmap
                    )
                except:
                    pass  # Shadow is optional, continue without it
                
                # Add the actual image with padding (matching web UI) and make it rounded
                image_shape = slide.shapes.add_picture(
                    image_path, 
                    Inches(0.55),  # left (with smaller padding)
                    Inches(1.55),  # top (with smaller padding) 
                    Inches(5.2),   # width reduced from 5.9 to 5.2
                    Inches(3.9)    # height (accounting for smaller padding)
                )
                
                # Apply rounded corners to the image itself
                try:
                    from pptx.oxml.xmlchemy import OxmlElement
                    from lxml import etree
                    pic_element = image_shape._element
                    pic_pr = pic_element.xpath("p:spPr")[0]
                    nsmap = pic_pr.nsmap
                    
                    # Add rounded rectangle geometry to image
                    prst_geom = etree.SubElement(pic_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 8333")  # Smaller rounded corner radius for image
                except:
                    pass  # Rounded image is optional
                
                # Add decorative circles to match web UI (positioned closer to rounded border edges)
                # Top-right circle (primary color) - positioned closer to rounded corner
                circle1 = slide.shapes.add_shape(
                    9,  # Oval shape
                    Inches(5.65), Inches(1.35), Inches(0.25), Inches(0.25)  # Adjusted for smaller image
                )
                circle1.fill.solid()
                circle1.fill.fore_color.rgb = slide_style["circle_colors"][0]
                circle1.line.fill.background()
                
                # Bottom-left circle (purple) - positioned closer to rounded corner
                circle2 = slide.shapes.add_shape(
                    9,  # Oval shape
                    Inches(0.35), Inches(5.25), Inches(0.2), Inches(0.2)  # Closer to rounded corner
                )
                circle2.fill.solid()
                circle2.fill.fore_color.rgb = slide_style["circle_colors"][1]
                
            except Exception as e:
                print(f"Error adding image to slide: {e}")
                
            # Text on the right with more gap
            text_left = Inches(6.7)  # Adjusted for smaller image (reduced from 7.5)
            text_width = Inches(6.1)  # Increased width (from 5.33)
        else:
            # Full width text if no image
            text_left = Inches(1)
            text_width = Inches(11.33)
        
        # Add content for image_left layout
        self._add_slide_content(slide, filtered_slide_data, slide_style, text_left, text_width)
    
    def _create_image_right_layout(self, slide, slide_data, image_path: Optional[str], slide_style: dict):
        """Image on right, text on left layout"""
        
        # Filter content to show only bullet points
        filtered_slide_data = self.FilteredSlideData(
            title=slide_data.title,
            content=self._filter_content_for_layout(slide_data.content, 'image_right'),
            image_url=slide_data.image_url,
            image_alt=slide_data.image_alt,
            layout=slide_data.layout
        )
        
        if image_path and os.path.exists(image_path):
            # Image on the right side
            try:
                # Add white background container (mimicking the web UI design)
                bg_container = slide.shapes.add_shape(
                    1,  # Rectangle shape
                    Inches(7.6), Inches(1.4), Inches(5.5), Inches(4.2)  # Right side position, reduced width from 6.2 to 5.5
                )
                bg_container.fill.solid()
                bg_container.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                bg_container.line.color.rgb = RGBColor(229, 231, 235)  # Light gray border
                bg_container.line.width = Pt(1)
                
                # Add rounded corners
                try:
                    from pptx.oxml.xmlchemy import OxmlElement
                    from lxml import etree
                    sp_element = bg_container._element
                    sp_pr = sp_element.xpath("p:spPr")[0]
                    nsmap = sp_pr.nsmap
                    
                    prst_geom = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 8333")
                except:
                    pass
                
                # Add the actual image
                image_shape = slide.shapes.add_picture(
                    image_path, 
                    Inches(7.75), Inches(1.55), Inches(5.2), Inches(3.9)  # Adjusted left position and reduced width from 5.9 to 5.2
                )
                
                # Apply rounded corners to image
                try:
                    pic_element = image_shape._element
                    pic_pr = pic_element.xpath("p:spPr")[0]
                    nsmap = pic_pr.nsmap
                    
                    prst_geom = etree.SubElement(pic_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 8333")
                except:
                    pass
                
                # Add decorative circles
                circle1 = slide.shapes.add_shape(
                    9, Inches(12.65), Inches(1.35), Inches(0.25), Inches(0.25)  # Adjusted for smaller image
                )
                circle1.fill.solid()
                circle1.fill.fore_color.rgb = slide_style["circle_colors"][0]
                circle1.line.fill.background()
                
                circle2 = slide.shapes.add_shape(
                    9, Inches(7.55), Inches(5.25), Inches(0.2), Inches(0.2)  # Adjusted for smaller image
                )
                circle2.fill.solid()
                circle2.fill.fore_color.rgb = slide_style["circle_colors"][1]
                circle2.line.fill.background()
                
            except Exception as e:
                print(f"Error adding image to slide: {e}")
            
            # Text on the left
            text_left = Inches(1)
            text_width = Inches(6.1)  # Increased from 5.5 to match smaller image
        else:
            # Full width text if no image
            text_left = Inches(1)
            text_width = Inches(11.33)
        
        self._add_slide_content(slide, filtered_slide_data, slide_style, text_left, text_width)
    
    def _create_text_only_layout(self, slide, slide_data, slide_style: dict):
        """Text-only layout, no images - centered content"""
        
        # Center everything
        text_left = Inches(2)
        text_width = Inches(9.33)
        
        # Add simple decorative elements
        # Small decorative circles
        circle1 = slide.shapes.add_shape(
            9, Inches(1), Inches(1), Inches(0.3), Inches(0.3)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = slide_style["accent_color"]
        circle1.line.fill.background()
        
        circle2 = slide.shapes.add_shape(
            9, Inches(11.5), Inches(6), Inches(0.4), Inches(0.4)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = slide_style["circle_colors"][1]
        circle2.line.fill.background()
        
        self._add_slide_content(slide, slide_data, slide_style, text_left, text_width)
    
    def _create_split_content_layout(self, slide, slide_data, image_path: Optional[str], slide_style: dict):
        """Split content layout - matching web UI 3-column design exactly"""
        
        # Add title spanning full width at top
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.33), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        title_font = title_para.font
        title_font.name = 'Inter'
        title_font.size = Pt(28)  # Smaller to match web UI
        title_font.bold = True
        title_font.color.rgb = slide_style.get('text_color', RGBColor(17, 24, 39))
        
        # Add accent line below title
        accent_line = slide.shapes.add_shape(
            1, Inches(6), Inches(1.3), Inches(1.33), Inches(0.06)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = slide_style["accent_color"]
        accent_line.line.fill.background()
        
        # Three column layout: bullets, image, description
        if image_path and os.path.exists(image_path):
            # Left Column - Bullets (matching web UI spacing)
            if slide_data.content:
                bullet_lines = [line.strip()[1:].strip() for line in slide_data.content.split('\n') 
                               if line.strip().startswith('-')][:4]  # Max 4 bullets
                
                bullets_box = slide.shapes.add_textbox(
                    Inches(0.3), Inches(2), Inches(3.8), Inches(4)
                )
                bullets_frame = bullets_box.text_frame
                bullets_frame.word_wrap = True
                bullets_frame.margin_left = Pt(0)
                bullets_frame.margin_right = Pt(0)
                bullets_frame.margin_top = Pt(0)
                bullets_frame.margin_bottom = Pt(0)
                
                for i, bullet_text in enumerate(bullet_lines):
                    if i == 0:
                        para = bullets_frame.paragraphs[0]
                    else:
                        para = bullets_frame.add_paragraph()
                    
                    # Add green dot bullet to match web UI
                    dot_run = para.add_run()
                    dot_run.text = "●"
                    dot_run.font.name = 'Inter'
                    dot_run.font.size = Pt(12)
                    dot_run.font.color.rgb = slide_style["accent_color"]
                    
                    # Add spacing
                    space_run = para.add_run()
                    space_run.text = "   "  # 3 spaces for proper alignment
                    
                    # Add text
                    text_run = para.add_run()
                    text_run.text = bullet_text
                    text_run.font.name = 'Inter'
                    text_run.font.size = Pt(13)  # Smaller to match web UI
                    text_run.font.color.rgb = slide_style.get('text_color', RGBColor(55, 65, 81))
                    
                    para.space_after = Pt(18)  # Larger spacing like web UI
                    para.alignment = PP_ALIGN.LEFT
            
            # Center Column - Image with rounded container (matching web UI)
            try:
                # Add rounded image container to match web UI
                img_container = slide.shapes.add_shape(
                    1, Inches(4.5), Inches(2), Inches(4.33), Inches(4)
                )
                img_container.fill.solid()
                img_container.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White like web UI
                img_container.line.color.rgb = RGBColor(209, 213, 219)  # Light border
                img_container.line.width = Pt(1)
                
                # Add rounded corners to container
                try:
                    from pptx.oxml.xmlchemy import OxmlElement
                    from lxml import etree
                    sp_element = img_container._element
                    sp_pr = sp_element.xpath("p:spPr")[0]
                    nsmap = sp_pr.nsmap
                    
                    prst_geom = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 8333")
                except:
                    pass
                
                # Add shadow effect
                try:
                    effect_list = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}effectLst", nsmap=nsmap)
                    outer_shadow = etree.SubElement(
                        effect_list, f"{{{nsmap['a']}}}outerShdw",
                        {"blurRad": "50800", "dist": "19050", "dir": "5400000", "algn": "tl"},
                        nsmap=nsmap
                    )
                    color_element = etree.SubElement(
                        outer_shadow, f"{{{nsmap['a']}}}srgbClr", {"val": "000000"}, nsmap=nsmap
                    )
                    etree.SubElement(
                        color_element, f"{{{nsmap['a']}}}alpha", {"val": "12000"}, nsmap=nsmap
                    )
                except:
                    pass
                
                # Add image with padding
                image_shape = slide.shapes.add_picture(
                    image_path, 
                    Inches(4.65), Inches(2.15), Inches(4.03), Inches(3.7)
                )
                
                # Round the image corners too
                try:
                    pic_element = image_shape._element
                    pic_pr = pic_element.xpath("p:spPr")[0]
                    nsmap = pic_pr.nsmap
                    
                    prst_geom = etree.SubElement(pic_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 8333")
                except:
                    pass
                    
            except Exception as e:
                print(f"Error adding image to slide: {e}")
            
            # Right Column - Description box (matching web UI styling)
            desc_box = slide.shapes.add_shape(
                1, Inches(9.2), Inches(2), Inches(3.6), Inches(4)
            )
            desc_box.fill.solid()
            desc_box.fill.fore_color.rgb = RGBColor(249, 250, 251)  # Light gray like web UI
            desc_box.line.color.rgb = RGBColor(229, 231, 235)
            desc_box.line.width = Pt(1)
            
            # Round the description box
            try:
                from pptx.oxml.xmlchemy import OxmlElement
                from lxml import etree
                sp_element = desc_box._element
                sp_pr = sp_element.xpath("p:spPr")[0]
                nsmap = sp_pr.nsmap
                
                prst_geom = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                prst_geom.set("prst", "roundRect")
                av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                gd.set("name", "adj")
                gd.set("fmla", "val 8333")
            except:
                pass
            
            # Add description text
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.margin_left = Pt(20)
            desc_frame.margin_right = Pt(20)
            desc_frame.margin_top = Pt(20)
            desc_frame.margin_bottom = Pt(20)
            
            # Add heading
            heading_para = desc_frame.paragraphs[0]
            heading_para.text = "Ключевые выводы"
            heading_para.alignment = PP_ALIGN.LEFT
            
            heading_font = heading_para.font
            heading_font.name = 'Inter'
            heading_font.size = Pt(16)
            heading_font.bold = True
            heading_font.color.rgb = RGBColor(17, 24, 39)
            
            # Add description text - use dynamic content or fallback
            desc_para = desc_frame.add_paragraph()
            
            # Extract key insights from content (last non-bullet line)
            if slide_data.content:
                # Get non-bullet lines that come after bullets (key insights)
                lines = slide_data.content.split('\n')
                non_bullet_lines = [line.strip() for line in lines 
                                   if line.strip() and not line.strip().startswith('-')]
                # Take the last non-empty line as key insights
                desc_text = non_bullet_lines[-1] if non_bullet_lines else 'Стратегические выводы и основные рекомендации по этому разделу.'
            else:
                desc_text = 'Основные выводы и рекомендации по данной теме.'
            
            desc_para.text = desc_text
            desc_para.alignment = PP_ALIGN.LEFT
            desc_para.space_before = Pt(12)
            
            desc_font = desc_para.font
            desc_font.name = 'Inter'
            desc_font.size = Pt(13)
            desc_font.color.rgb = RGBColor(75, 85, 99)
            
            # Add decorative circles (matching web UI)
            circles_para = desc_frame.add_paragraph()
            circles_para.text = "● ● ●"
            circles_para.alignment = PP_ALIGN.LEFT
            circles_para.space_before = Pt(16)
            
            circles_font = circles_para.font
            circles_font.name = 'Inter'
            circles_font.size = Pt(12)
            circles_font.color.rgb = slide_style["accent_color"]
            
        else:
            # Fall back to regular layout
            self._add_slide_content(slide, slide_data, slide_style, Inches(1), Inches(11.33))
    
    def _create_image_top_layout(self, slide, slide_data, image_path: Optional[str], slide_style: dict):
        """Image on top, text below layout"""
        
        # Filter content to show only bullet points
        filtered_slide_data = self.FilteredSlideData(
            title=slide_data.title,
            content=self._filter_content_for_layout(slide_data.content, 'image_top'),
            image_url=slide_data.image_url,
            image_alt=slide_data.image_alt,
            layout=slide_data.layout
        )
        
        if image_path and os.path.exists(image_path):
            # Image at top - full width with cover behavior to match web UI
            try:
                # Get image dimensions
                from PIL import Image as PILImage
                with PILImage.open(image_path) as img:
                    img_width, img_height = img.size
                    img_aspect = img_width / img_height
                
                # Target area dimensions
                target_width = 13.33  # inches
                target_height = 3.2   # inches
                target_aspect = target_width / target_height
                
                # Add picture to slide
                pic = slide.shapes.add_picture(
                    image_path, 
                    Inches(0), Inches(0), Inches(target_width), Inches(target_height)
                )
                
                # Apply cropping to achieve cover behavior
                # Crop percentages are in ratio (0.0 to 1.0 represents 0% to 100%)
                if img_aspect > target_aspect:
                    # Image is wider - need to crop left and right
                    # Calculate how much wider the image is
                    scale_factor = target_height / (img_height / 96)  # 96 DPI assumption
                    scaled_width = (img_width / 96) * scale_factor
                    crop_total = (scaled_width - target_width) / scaled_width
                    
                    # Crop equally from left and right
                    pic.crop_left = crop_total / 2
                    pic.crop_right = crop_total / 2
                else:
                    # Image is taller - need to crop top and bottom
                    scale_factor = target_width / (img_width / 96)  # 96 DPI assumption
                    scaled_height = (img_height / 96) * scale_factor
                    crop_total = (scaled_height - target_height) / scaled_height
                    
                    # Crop equally from top and bottom
                    pic.crop_top = crop_total / 2
                    pic.crop_bottom = crop_total / 2
                    
            except Exception as e:
                print(f"Error adding image to slide: {e}")
            
            # Text below image - adjusted for larger image
            text_left = Inches(1)
            text_top = Inches(3.4)  # Reduced from 3.8 to 3.4 to prevent overflow
            text_width = Inches(11.33)
        else:
            # No image, normal text positioning
            text_left = Inches(1)
            text_top = Inches(1)
            text_width = Inches(11.33)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            text_left, text_top, text_width, Inches(0.8)  # Reduced from 1 to 0.8
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER if slide_style["text_alignment"] == "center" else PP_ALIGN.LEFT
        
        title_font = title_para.font
        title_font.name = 'Inter'
        title_font.size = Pt(32)  # Reduced from 36 to 32
        title_font.bold = True
        title_font.color.rgb = slide_style.get('text_color', RGBColor(17, 24, 39))
        
        # Add accent line
        line_left = text_left if slide_style["text_alignment"] == "left" else Inches(6)
        accent_line = slide.shapes.add_shape(
            1, line_left, text_top + Inches(0.9), Inches(1.2), Inches(0.06)  # Reduced from 1.1 to 0.9
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = slide_style["accent_color"]
        accent_line.line.fill.background()
        
        # Add content
        if filtered_slide_data.content:
            content_box = slide.shapes.add_textbox(
                text_left, text_top + Inches(1.2), text_width, Inches(3.3)  # Reduced start from 1.5 to 1.2, increased height from 3 to 3.3
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            self._process_bullet_content(content_frame, filtered_slide_data.content, slide_style)
    
    def _create_grid_layout(self, slide, slide_data, image_path: Optional[str], slide_style: dict):
        """Grid layout - content in structured boxes"""
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.33), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        title_font = title_para.font
        title_font.name = 'Inter'
        title_font.size = Pt(36)
        title_font.bold = True
        title_font.color.rgb = slide_style.get('text_color', RGBColor(17, 24, 39))
        
        # Add accent line
        accent_line = slide.shapes.add_shape(
            1, Inches(5.5), Inches(1.6), Inches(2.33), Inches(0.06)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = slide_style["accent_color"]
        accent_line.line.fill.background()
        
        if slide_data.content:
            lines = [line.strip() for line in slide_data.content.split('\n') if line.strip().startswith('-')]
            
            # Create grid boxes for content
            box_width = Inches(5.5)
            box_height = Inches(2.2)
            
            positions = [
                (Inches(1), Inches(2.3)),      # Top left
                (Inches(7.3), Inches(2.3)),    # Top right  
                (Inches(1), Inches(4.8)),      # Bottom left
                (Inches(7.3), Inches(4.8))     # Bottom right
            ]
            
            for i, line in enumerate(lines[:4]):  # Max 4 boxes
                if i >= len(positions):
                    break
                    
                pos_x, pos_y = positions[i]
                
                # Create box
                box = slide.shapes.add_shape(
                    1, pos_x, pos_y, box_width, box_height
                )
                box.fill.solid()
                # Use dark gray for dark theme, light gray for others
                box_bg = RGBColor(31, 41, 55) if slide_style.get('content_bg') == RGBColor(17, 24, 39) else RGBColor(248, 250, 252)
                box.fill.fore_color.rgb = box_bg
                box.line.color.rgb = slide_style["accent_color"]
                box.line.width = Pt(2)
                
                # Add rounded corners
                try:
                    from pptx.oxml.xmlchemy import OxmlElement
                    from lxml import etree
                    sp_element = box._element
                    sp_pr = sp_element.xpath("p:spPr")[0]
                    nsmap = sp_pr.nsmap
                    
                    prst_geom = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}prstGeom", nsmap=nsmap)
                    prst_geom.set("prst", "roundRect")
                    av_lst = etree.SubElement(prst_geom, f"{{{nsmap['a']}}}avLst", nsmap=nsmap)
                    gd = etree.SubElement(av_lst, f"{{{nsmap['a']}}}gd", nsmap=nsmap)
                    gd.set("name", "adj")
                    gd.set("fmla", "val 16667")  # More rounded
                except:
                    pass
                
                # Add number badge in corner
                badge = slide.shapes.add_shape(
                    9, pos_x + Inches(0.3), pos_y + Inches(0.3), Inches(0.5), Inches(0.5)
                )
                badge.fill.solid()
                badge.fill.fore_color.rgb = slide_style["accent_color"]
                badge.line.fill.background()
                
                # Add number to badge
                badge_frame = badge.text_frame
                badge_para = badge_frame.paragraphs[0]
                badge_para.text = str(i + 1)
                badge_para.alignment = PP_ALIGN.CENTER
                badge_font = badge_para.font
                badge_font.name = 'Inter'
                badge_font.size = Pt(14)
                badge_font.bold = True
                badge_font.color.rgb = RGBColor(255, 255, 255)
                
                # Add text to box
                text_frame = box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Pt(24)
                text_frame.margin_right = Pt(24)
                text_frame.margin_top = Pt(40)  # Leave space for badge
                text_frame.margin_bottom = Pt(24)
                
                para = text_frame.paragraphs[0]
                para.text = line.strip()[1:].strip()  # Remove the dash
                para.alignment = PP_ALIGN.CENTER
                
                font = para.font
                font.name = 'Inter'
                font.size = Pt(16)
                font.color.rgb = slide_style.get('text_color', RGBColor(55, 65, 81))
    
    def _add_slide_content(self, slide, slide_data, slide_style: dict, text_left, text_width):
        """Add title, accent line, and content to slide"""
        
        # Add title
        title_box = slide.shapes.add_textbox(
            text_left, Inches(1), text_width, Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        
        # Apply alignment based on style
        if slide_style["text_alignment"] == "center":
            title_para.alignment = PP_ALIGN.CENTER
        else:
            title_para.alignment = PP_ALIGN.LEFT
        
        # Title font
        title_font = title_para.font
        title_font.name = 'Inter'
        title_font.size = Pt(36)
        title_font.bold = True
        title_font.color.rgb = slide_style.get('text_color', RGBColor(17, 24, 39))
        
        # Add accent line (always present, matching web UI)
        line_left = text_left if slide_style["text_alignment"] == "left" else text_left + (text_width / 2) - Inches(0.6)
        accent_line = slide.shapes.add_shape(
            1,  # Rectangle shape
            line_left, Inches(2.6), Inches(1.2), Inches(0.06)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = slide_style["accent_color"]
        accent_line.line.fill.background()
        
        # Add content with proper bullet points
        if slide_data.content:
            content_box = slide.shapes.add_textbox(
                text_left, Inches(3), text_width, Inches(4)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            self._process_bullet_content(content_frame, slide_data.content, slide_style)
    
    def _process_bullet_content(self, content_frame, content_text: str, slide_style: dict):
        """Process bullet point content with style variations"""
        
        content_lines = content_text.split('\n')
        
        first_para = True
        bullet_count = 0
        
        for line in content_lines:
            if not line.strip():
                continue
                
            if first_para:
                para = content_frame.paragraphs[0]
                first_para = False
            else:
                para = content_frame.add_paragraph()
            
            # Handle bullet points with style variations
            if line.strip().startswith('-'):
                bullet_count += 1
                
                if slide_style["bullet_type"] == "numbers":
                    # Create numbered list (01, 02, 03, etc.)
                    bullet_run = para.add_run()
                    bullet_run.text = f"{bullet_count:02d}"  # Format as 01, 02, 03...
                    bullet_run.font.name = 'Inter'
                    bullet_run.font.size = Pt(16)  # Slightly smaller for numbers
                    bullet_run.font.bold = True
                    bullet_run.font.color.rgb = slide_style["accent_color"]
                    
                    # Add period and spacing
                    period_run = para.add_run()
                    period_run.text = ".  "  # Period and spaces
                    period_run.font.name = 'Inter'
                    period_run.font.size = Pt(16)
                    period_run.font.color.rgb = slide_style["accent_color"]
                else:
                    # Create bullet dot
                    bullet_run = para.add_run()
                    bullet_run.text = "●"  # Just the bullet character
                    bullet_run.font.name = 'Inter'
                    bullet_run.font.size = Pt(20)  # Larger size for visibility
                    bullet_run.font.color.rgb = slide_style["accent_color"]
                    
                    # Add spacing after bullet
                    space_run = para.add_run()
                    space_run.text = "  "  # Two spaces for proper spacing
                
                # Add text content
                text_run = para.add_run()
                text_run.text = line.strip()[1:].strip()
                text_run.font.name = 'Inter'
                text_run.font.size = Pt(18)  # text-lg in web UI
                text_run.font.color.rgb = slide_style.get('text_color', RGBColor(55, 65, 81))
                
                # Set paragraph alignment
                para.alignment = PP_ALIGN.LEFT
                
                # Set spacing between bullets to match mb-4 (16px)
                para.space_after = Pt(16)  # Matching mb-4
                para.space_before = Pt(0)  # No space before
                
                # Adjust line spacing for better bullet alignment
                para.line_spacing = 1.2  # Tighter line spacing
                
            else:
                # Non-bullet text
                para.text = self._clean_markdown(line)
                font = para.font
                font.name = 'Inter'
                font.size = Pt(18)  # text-lg to match web UI
                font.color.rgb = slide_style.get('text_color', RGBColor(55, 65, 81))
    
    def _clean_markdown(self, text: str) -> str:
        """Clean markdown formatting for PPTX"""
        if not text:
            return ""
        
        # Remove markdown formatting
        text = text.replace('**', '')
        text = text.replace('*', '')
        text = text.replace('###', '')
        text = text.replace('##', '')
        text = text.replace('#', '')
        
        return text.strip()