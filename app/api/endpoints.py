from fastapi import APIRouter, Depends, HTTPException
from sqlalchemy.orm import Session
from .. import models, schemas, database
import httpx, os, json, secrets, tempfile, re
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any, Optional
from urllib.parse import urlparse
from typing import Tuple

router = APIRouter()


def read_system_prompt() -> str:
    with open("prompts/system.txt", "r", encoding="utf-8") as f:
        return f.read().strip()


def generate_presentation_uid() -> str:
    token = secrets.token_urlsafe(8)
    return f"pres_{token}"


def _nice_title(s: str) -> str:
    if not s:
        return ""
    return re.sub(r"[_\-]+", " ", s).strip().title()


@router.post("/templates/", response_model=schemas.PresentationTemplate)
def create_template(template: schemas.PresentationTemplateCreate, db: Session = Depends(database.get_db)):
    db_template = models.PresentationTemplate(**template.model_dump())
    db.add(db_template)
    db.commit()
    db.refresh(db_template)
    return db_template


@router.get("/templates/", response_model=list[schemas.PresentationTemplate])
def read_templates(skip: int = 0, limit: int = 100, db: Session = Depends(database.get_db)):
    return db.query(models.PresentationTemplate).offset(skip).limit(limit).all()


@router.post("/presentations/", response_model=schemas.Presentation)
def create_presentation(presentation: schemas.PresentationCreate, db: Session = Depends(database.get_db)):
    db_pres = models.Presentation(**presentation.model_dump())
    db.add(db_pres)
    db.commit()
    db.refresh(db_pres)
    return db_pres


@router.get("/presentations/", response_model=list[schemas.Presentation])
def read_presentations(skip: int = 0, limit: int = 100, db: Session = Depends(database.get_db)):
    return db.query(models.Presentation).offset(skip).limit(limit).all()


async def fetch_tumblr_images(query: str, limit: int = 1, before: Optional[int] = None) -> List[Dict[str, Any]]:
    api_key = os.getenv("TUMBLR_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="TUMBLR_API_KEY not set")

    limit = max(1, min(int(limit or 1), 20))
    fetch_limit = min(max(limit * 5, limit), 20)

    params = {
        "tag": query,
        "api_key": api_key,
        "limit": fetch_limit,
        "npf": "true",
    }
    if before:
        params["before"] = int(before)

    async with httpx.AsyncClient(timeout=30.0) as client:
        r = await client.get("https://api.tumblr.com/v2/tagged", params=params)
        try:
            r.raise_for_status()
        except httpx.HTTPStatusError as e:
            raise HTTPException(status_code=502, detail=f"Tumblr error: {e.response.text}")

        items = (r.json().get("response") or [])

    def _is_mature(post: Dict[str, Any]) -> bool:
        rating = str(post.get("content_rating") or "").lower()
        if rating in ("mature", "adult"):
            return True
        if post.get("is_nsfw") or post.get("is_adult"):
            return True
        blog = post.get("blog") or {}
        if isinstance(blog, dict) and (blog.get("is_nsfw") or blog.get("is_adult")):
            return True
        return False

    def _largest_media(media: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
        if not media:
            return None
        return max(media, key=lambda m: (m.get("width") or 0, m.get("height") or 0))

    def _extract_images(post: Dict[str, Any]) -> List[Tuple[str, Optional[int], Optional[int]]]:
        out: List[Tuple[str, Optional[int], Optional[int]]] = []

        photos = post.get("photos")
        if isinstance(photos, list):
            for p in photos:
                size = p.get("original_size")
                if not size:
                    sizes = p.get("alt_sizes") or []
                    if sizes:
                        size = max(sizes, key=lambda s: s.get("width", 0))
                if size and size.get("url"):
                    out.append((size["url"], size.get("width"), size.get("height")))

        content = post.get("content")
        if isinstance(content, list):
            for block in content:
                if block.get("type") == "image":
                    media = block.get("media") or []
                    best = _largest_media(media)
                    if best and best.get("url"):
                        out.append((best["url"], best.get("width"), best.get("height")))
        return out

    results: List[Dict[str, Any]] = []
    for item in items:
        if _is_mature(item):
            continue
        title = item.get("summary") or item.get("blog_name")
        post_url = item.get("post_url")
        for url, w, h in _extract_images(item):
            results.append({
                "url": url,
                "width": w, "height": h,
                "title": title,
                "post_url": post_url,
                "source": "tumblr",
                "query": query,
            })
            if len(results) >= limit:
                break
        if len(results) >= limit:
            break

    return results[:limit]


async def enrich_presentation_with_images(pres: Dict[str, Any], default_query: Optional[str] = None) -> Dict[str, Any]:
    slides = pres.get("slides", []) or []
    for slide in slides:
        fields = slide.get("fields")
        if isinstance(fields, dict):
            for key, fld in list(fields.items()):
                val = (fld or {}).get("value")
                is_image_block = isinstance(val, dict) and val.get("type") == "image"
                looks_like_image = isinstance(val, dict) and ("query" in val or "url" in val)
                if is_image_block or looks_like_image:
                    content = val.get("content", val) if isinstance(val, dict) else val
                    query = content.get("query") if isinstance(content, dict) else None
                    query = query or default_query
                    if query and not content.get("url"):
                        imgs = await fetch_tumblr_images(query, limit=1)
                        if imgs:
                            if is_image_block:
                                val["content"] = imgs[0]
                                fld["value"] = val
                            else:
                                fld["value"] = imgs[0]
                            fields[key] = fld
        elif isinstance(fields, list):
            for fld in fields:
                if not isinstance(fld, dict):
                    continue
                if fld.get("type") == "image":
                    v = fld.get("value")
                    query = v.get("query") if isinstance(v, dict) else (v if isinstance(v, str) else None)
                    if query:
                        imgs = await fetch_tumblr_images(query, limit=1)
                        if imgs:
                            fld["value"] = imgs[0]
    pres["slides"] = slides
    return pres


@router.post("/generate", response_model=schemas.Presentation)
async def generate_presentation(request: schemas.GenerateRequest, db: Session = Depends(database.get_db)):
    template = db.query(models.PresentationTemplate).filter(
        models.PresentationTemplate.id == request.template_id
    ).first()
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")

    tpl = template.slides
    if isinstance(tpl, dict) and "slides" in tpl:
        slide_template_obj = tpl
    else:
        slide_template_obj = {
            "id": getattr(template, "slug", None) or f"tpl_{template.id}",
            "title": getattr(template, "title", None) or (request.title or "Presentation"),
            "slides": tpl if isinstance(tpl, list) else []
        }

    system_prompt = read_system_prompt()

    user_prompt = (
        f'topic: "{request.title or request.description or "Presentation"}"\n\n'
        "slide_template:\n"
        f"{json.dumps(slide_template_obj, ensure_ascii=False, indent=2)}"
    )

    api_key = os.getenv("OPENROUTER_API_KEY")
    model = os.getenv("OPENROUTER_MODEL", "mistralai/mixtral-8x7b-instruct")
    if not api_key:
        raise HTTPException(status_code=500, detail="OPENROUTER_API_KEY not set")

    headers = {"Authorization": f"Bearer {api_key}", "HTTP-Referer": "http://localhost:8000",
               "X-Title": "Perio Backend"}
    payload = {"model": model, "messages": [{"role": "system", "content": system_prompt},
                                            {"role": "user", "content": user_prompt}], "temperature": 0.7}

    async with httpx.AsyncClient(timeout=60.0) as client:
        try:
            resp = await client.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload)
            resp.raise_for_status()
            data = resp.json()
            raw = data["choices"][0]["message"]["content"].strip()
            if raw.startswith("```"):
                parts = raw.split("```", 2)
                if len(parts) >= 2:
                    raw = parts[1]
                    if raw.startswith("json"):
                        raw = raw[4:].strip()

            presentation_data = json.loads(raw)

            presentation_data = await enrich_presentation_with_images(presentation_data, default_query=request.title)

            slides_count = len(presentation_data.get("slides", []))
            uid = generate_presentation_uid()
            base_pr_url = os.getenv("PRESENTATION_BASE_URL", "http://localhost:8000/presentations")
            full_url = f"{base_pr_url}/{uid}"

            db_presentation = models.Presentation(
                title=request.title,
                description=request.description,
                presentation=presentation_data,
                generating=True,
                slides_count=slides_count,
                presentation_template_id=request.template_id,
                presentation_url=full_url
            )
            db.add(db_presentation)
            db.commit()
            db.refresh(db_presentation)
            return db_presentation

        except httpx.HTTPStatusError as e:
            raise HTTPException(status_code=502, detail=f"OpenRouter error: {e.response.text}")
        except json.JSONDecodeError:
            raise HTTPException(status_code=502, detail="Invalid JSON from model")
        except Exception as e:
            db.rollback()
            raise HTTPException(status_code=500, detail=f"Generation failed: {str(e)}")


def _download_image_to_tmp(url: str) -> Optional[str]:
    try:
        with httpx.Client(timeout=30.0) as client:
            r = client.get(url)
            r.raise_for_status()
            suffix = ".png"
            path = urlparse(url).path.lower()
            for ext in (".png", ".jpg", ".jpeg", ".webp"):
                if path.endswith(ext):
                    suffix = ext
                    break
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as img_tmp:
                img_tmp.write(r.content)
                return img_tmp.name
    except Exception:
        return None


def _add_image_slide(prs: Presentation, image_path: str, caption: Optional[str] = None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    left, top = Inches(0.5), Inches(0.5)
    max_width = prs.slide_width - Inches(1.0)
    try:
        slide.shapes.add_picture(image_path, left, top, width=max_width)
    except Exception:
        slide.shapes.add_picture(image_path, left, top)
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(1.2), prs.slide_width - Inches(1.0),
                                      Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        p.text = str(caption)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)


def _flatten_list(v: Any) -> List[str]:
    out: List[str] = []
    if isinstance(v, list):
        for it in v:
            if isinstance(it, str):
                out.append(it)
            elif isinstance(it, dict):
                for val in it.values():
                    if isinstance(val, (str, int, float, bool)):
                        out.append(str(val))
                        break
    return out


def build_presentation_from_json(presentation_data: dict, title: str = "Presentation") -> str:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    slides: List[Dict[str, Any]] = presentation_data.get("slides", []) or []

    for slide_obj in slides:
        fields = slide_obj.get("fields", {})
        slide_title = slide_obj.get("title")

        if not isinstance(fields, dict):
            if isinstance(fields, list):
                fields = {f.get("id"): {"value": f.get("value")} for f in fields if isinstance(f, dict) and "id" in f}
            else:
                fields = {}

        fval = lambda k: (fields.get(k) or {}).get("value")

        if fval("title") or fval("subtitle"):
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title and fval("title"):
                slide.shapes.title.text = str(fval("title"))
            sub_parts = [str(x) for x in [fval("subtitle"), fval("presenter"), fval("date")] if x]
            if len(slide.placeholders) > 1 and sub_parts:
                slide.placeholders[1].text = "\n".join(sub_parts)
            for k, v in fields.items():
                val = v.get("value")
                if isinstance(val, dict) and val.get("url"):
                    img = _download_image_to_tmp(val["url"])
                    if img: _add_image_slide(prs, img, caption=val.get("title"))
            continue

        header = slide_title or fval("section") or fval("header") or _nice_title(slide_obj.get("id", "")) or title
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = str(header)

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        text_fields = []
        for k, v in fields.items():
            val = v.get("value")
            if isinstance(val, str) and k not in ("title", "subtitle"):
                text_fields.append(val)
        for i, t in enumerate(text_fields):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = t
            p.level = 0

        bullets: List[str] = []
        for k, v in fields.items():
            val = v.get("value")
            if isinstance(val, list):
                bullets += _flatten_list(val)

        if bullets and not text_fields:
            tf.clear()
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if (text_fields or i > 0) else tf.paragraphs[0]
            p.text = b
            p.level = 0

        for k, v in fields.items():
            val = v.get("value")
            if isinstance(val, dict) and val.get("url"):
                img = _download_image_to_tmp(val["url"])
                if img:
                    _add_image_slide(prs, img, caption=val.get("title"))

    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = title

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        return tmp.name


@router.get("/presentations/{presentation_id}/download")
async def download_presentation(presentation_id: int, db: Session = Depends(database.get_db)):
    db_pres = db.query(models.Presentation).filter(models.Presentation.id == presentation_id).first()
    if not db_pres:
        raise HTTPException(status_code=404, detail="Presentation not found")
    if not db_pres.presentation:
        raise HTTPException(status_code=400, detail="Presentation data is empty")

    try:
        pptx_path = build_presentation_from_json(db_pres.presentation, title=db_pres.title or "Presentation")
        filename = f"presentation_{presentation_id}.pptx"
        return FileResponse(path=pptx_path, filename=filename,
                            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX: {str(e)}")


@router.get("/images/search")
async def search_images(query: str, limit: int = 1):
    imgs = await fetch_tumblr_images(query=query, limit=limit)
    return {"query": query, "images": imgs}
